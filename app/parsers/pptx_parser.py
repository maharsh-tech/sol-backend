from pptx import Presentation


def parse_pptx(file_path: str, document_name: str) -> list[dict]:
    """Parse a PowerPoint file and return a list of text chunks with metadata."""
    results: list[dict] = []
    prs = Presentation(file_path)
    for slide_num, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = paragraph.text.strip()
                    if para_text:
                        texts.append(para_text)
        text = "\n".join(texts).strip()
        if text:
            results.append(
                {
                    "text": text,
                    "metadata": {
                        "document_name": document_name,
                        "page": slide_num,
                        "source_type": "pptx",
                    },
                }
            )
    return results
